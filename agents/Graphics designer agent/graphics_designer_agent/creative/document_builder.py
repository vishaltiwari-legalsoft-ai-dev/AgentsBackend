"""Structured layout engine — render a reviewed plan into real, branded assets.

This is the "match or exceed the reference" engine. Unlike the social pipeline
(flat images from the image model), brochures and decks are built as *structured*
documents with live text and the brand's real palette + fonts:

- ``build_brochure_pdf``     — reportlab (selectable text, print-grade).
- ``build_presentation_pptx``— python-pptx (real, editable slides + speaker notes).
- ``build_carousel_frames``  — Pillow PNGs, one per frame, one consistent system.
- ``build_blog_images``      — Pillow cover + in-article PNGs.

Heavy export deps are imported lazily so the package still imports without them;
``engine_status()`` reports which engines are actually available on this host.
Each builder resolves the brand identity from a ``BrandPack`` (palette + fonts +
name), so output is on-brand for whichever brand the run selected.
"""

from __future__ import annotations

import io
import logging
import zipfile
from typing import Any, Callable, Optional

from PIL import Image, ImageDraw, ImageFont

from . import brochure_layout, brochure_render

logger = logging.getLogger("graphics_designer.creative.document_builder")

# An artifact is (filename, bytes, mime-type).
Artifact = tuple[str, bytes, str]
# Optional per-artifact callback — invoked the moment each artifact is ready so the
# caller can persist it + advance a progress bar while a multi-file build streams in.
OnArtifact = Optional[Callable[[Artifact], None]]


def _emit(on_artifact: OnArtifact, artifacts: list[Artifact]) -> list[Artifact]:
    """Fire ``on_artifact`` for each item (used by single-shot builders that don't
    stream). Best-effort — a failing callback never breaks the build."""
    if on_artifact:
        for art in artifacts:
            try:
                on_artifact(art)
            except Exception:  # noqa: BLE001
                logger.warning("on_artifact callback failed", exc_info=True)
    return artifacts


# --------------------------------------------------------------------------- #
# Brand identity resolution (defensive — works for any BrandPack shape)
# --------------------------------------------------------------------------- #

def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    v = (value or "").lstrip("#")
    if len(v) == 3:
        v = "".join(c * 2 for c in v)
    if len(v) != 6:
        return (15, 15, 15)
    try:
        return (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
    except ValueError:
        return (15, 15, 15)


def _brand_palette(pack: Any) -> dict[str, Any]:
    lc = getattr(pack, "locked_colors", {}) or {}
    grad = list(lc.get("gradient") or ["#FFFFFF", "#1746A2"])
    if len(grad) < 2:
        grad = ["#FFFFFF", grad[0] if grad else "#1746A2"]
    cta = lc.get("cta") if isinstance(lc.get("cta"), dict) else {}
    deep = grad[-1]
    return {
        "gradient": grad,
        "light": grad[0],
        "deep": deep,
        "accent": lc.get("accent", deep),
        "text": lc.get("text", "#0F0F0F"),
        "cta_from": cta.get("from", lc.get("accent", deep)),
        "cta_to": cta.get("to", deep),
    }


def _font_path(pack: Any, name: Optional[str] = None) -> Optional[str]:
    """Absolute path to a brand font file, or None if it can't be resolved."""
    try:
        fonts_dir = getattr(pack, "fonts_dir", None)
        if not fonts_dir:
            return None
        fname = pack.font_file(name or getattr(pack, "default_font", ""))
        path = fonts_dir / fname
        return str(path) if path.exists() else None
    except Exception:  # noqa: BLE001
        return None


def _load_font(pack: Any, size: int, name: Optional[str] = None) -> ImageFont.FreeTypeFont:
    path = _font_path(pack, name)
    if path:
        try:
            return ImageFont.truetype(path, size=size)
        except Exception:  # noqa: BLE001
            pass
    try:
        return ImageFont.truetype("arial.ttf", size=size)
    except Exception:  # noqa: BLE001
        return ImageFont.load_default()


# --------------------------------------------------------------------------- #
# Pillow helpers (carousel + blog raster output)
# --------------------------------------------------------------------------- #

def _vertical_gradient(size: tuple[int, int], top: tuple[int, int, int],
                       bottom: tuple[int, int, int]) -> Image.Image:
    w, h = size
    base = Image.new("RGB", size, top)
    draw = ImageDraw.Draw(base)
    for y in range(h):
        t = y / max(1, h - 1)
        r = int(top[0] + (bottom[0] - top[0]) * t)
        g = int(top[1] + (bottom[1] - top[1]) * t)
        b = int(top[2] + (bottom[2] - top[2]) * t)
        draw.line([(0, y), (w, y)], fill=(r, g, b))
    return base


def _wrap(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont,
          max_w: int) -> list[str]:
    words = (text or "").split()
    lines: list[str] = []
    cur = ""
    for w in words:
        trial = f"{cur} {w}".strip()
        if draw.textlength(trial, font=font) <= max_w or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


def _draw_block(draw: ImageDraw.ImageDraw, x: int, y: int, lines: list[str],
                font: ImageFont.FreeTypeFont, fill: tuple[int, int, int],
                line_gap: float = 1.25, center_w: Optional[int] = None) -> int:
    asc, desc = font.getmetrics()
    lh = int((asc + desc) * line_gap)
    for ln in lines:
        lx = x
        if center_w is not None:
            lx = x + (center_w - int(draw.textlength(ln, font=font))) // 2
        draw.text((lx, y), ln, font=font, fill=fill)
        y += lh
    return y


def _render_text_frame(
    pack: Any, size: tuple[int, int], *, headline: str, body: str = "",
    role: str = "body", footer: str = "", cta: str = "",
) -> bytes:
    """One branded raster frame: gradient bg, headline, body, optional CTA pill."""
    pal = _brand_palette(pack)
    w, h = size
    light = _hex_to_rgb(pal["light"])
    deep = _hex_to_rgb(pal["deep"])
    # Hook/cover lead light→deep; quieter frames stay light for legible body copy.
    top, bottom = (light, deep) if role in ("hook", "cover", "cta") else (
        (255, 255, 255), light)
    img = _vertical_gradient(size, top, bottom)
    draw = ImageDraw.Draw(img)
    on_dark = role in ("hook", "cover", "cta")
    ink = (255, 255, 255) if on_dark else _hex_to_rgb(pal["text"])

    margin = int(w * 0.09)
    max_w = w - 2 * margin
    head_font = _load_font(pack, int(h * 0.085), getattr(pack, "default_font", None))
    body_font = _load_font(pack, int(h * 0.038), getattr(pack, "default_font", None))

    head_lines = _wrap(draw, headline, head_font, max_w)
    body_lines = _wrap(draw, body, body_font, max_w) if body else []
    asc, desc = head_font.getmetrics()
    head_h = int((asc + desc) * 1.18) * len(head_lines)
    basc, bdesc = body_font.getmetrics()
    body_h = int((basc + bdesc) * 1.3) * len(body_lines)
    total = head_h + (int(h * 0.03) + body_h if body_lines else 0)
    y = max(margin, (h - total) // 2)

    y = _draw_block(draw, margin, y, head_lines, head_font, ink, 1.18)
    if body_lines:
        y += int(h * 0.03)
        body_ink = (235, 240, 250) if on_dark else _hex_to_rgb(pal["deep"])
        y = _draw_block(draw, margin, y, body_lines, body_font, body_ink, 1.3)

    if cta:
        _draw_cta_pill(draw, img, pack, cta, w, h)

    if footer:
        f_font = _load_font(pack, int(h * 0.026))
        draw.text((margin, h - int(h * 0.06)), footer, font=f_font,
                  fill=(255, 255, 255) if on_dark else _hex_to_rgb(pal["deep"]))

    out = io.BytesIO()
    img.save(out, format="PNG")
    return out.getvalue()


def _draw_cta_pill(draw: ImageDraw.ImageDraw, img: Image.Image, pack: Any,
                   text: str, w: int, h: int) -> None:
    pal = _brand_palette(pack)
    font = _load_font(pack, int(h * 0.04), getattr(pack, "default_font", None))
    tw = int(draw.textlength(text, font=font))
    pad_x, pad_y = int(w * 0.05), int(h * 0.022)
    pill_w, pill_h = tw + 2 * pad_x, font.getmetrics()[0] + font.getmetrics()[1] + 2 * pad_y
    x0 = (w - pill_w) // 2
    y0 = int(h * 0.74)
    draw.rounded_rectangle([x0, y0, x0 + pill_w, y0 + pill_h], radius=pill_h // 2,
                           fill=_hex_to_rgb(pal["cta_from"]))
    draw.text((x0 + pad_x, y0 + pad_y), text, font=font, fill=(255, 255, 255))


# --------------------------------------------------------------------------- #
# Carousel + blog (Pillow) — always available
# --------------------------------------------------------------------------- #

# A safety ceiling on per-slide generations so a runaway count can't trigger a
# huge number of image-model calls. Normal carousels are 3–7 slides.
_MAX_GENERATED_FRAMES = 10
# Slides are independent, so generate them concurrently instead of back-to-back —
# wall-clock becomes ~one slide's latency instead of the sum of all of them. Capped
# to stay friendly to the image provider's rate limits.
_CAROUSEL_CONCURRENCY = 5


def _carousel_text_fallback(pack: Any, fr: dict, idx: int,
                            images_only: bool = False) -> Artifact:
    """Local Pillow gradient frame — used when a slide's image generation fails so
    the user still gets a downloadable, on-brand frame instead of nothing. In
    ``images_only`` mode the frame carries no copy (just the brand gradient), to
    match the "images with logo only" choice."""
    from .. import reference_library as rl
    w, h = rl.type_spec("carousel").get("target_dims", (1080, 1080))
    role = fr.get("role", "body")
    if images_only:
        png = _render_text_frame(pack, (w, h), headline="", body="", role="cover")
    else:
        png = _render_text_frame(
            pack, (w, h), headline=fr.get("headline", ""), body=fr.get("body", ""),
            role=role, footer=f"{getattr(pack, 'name', 'Brand')} · {idx}",
            cta=fr.get("headline", "") if role == "cta" else "",
        )
    return (f"frame-{idx:02d}.png", png, "image/png")


def build_carousel_frames(plan: dict[str, Any], pack: Any,
                          on_artifact: OnArtifact = None) -> list[Artifact]:
    """Render a TRUE carousel: a DISTINCT on-brand image per slide, not one photo ×N.

    Each slide runs the backbone's Stage 1 (brand foundation) + Stage 2 with its own
    ``subject`` from the plan, so the foregrounds differ while the shared brand
    palette/gradient/logo keep the set cohesive ("shared style, different subjects").
    Stage 3 (text, with the key phrase highlighted) + Stage 4 (logo) are then applied
    per slide via the same renderer the interactive editor uses, so text stays crisp,
    correct, and in the brand font.

    Slides are generated CONCURRENTLY (each is independent) so a 5-slide carousel
    takes roughly one slide's wall-clock, not five. A slide whose image generation
    fails individually falls back to a local gradient text-frame, so one bad slide
    never sinks the whole set.
    """
    import concurrent.futures

    from .. import pipeline, reference_library as rl
    from . import layout_brain

    frames = (plan.get("frames", []) or [])[:_MAX_GENERATED_FRAMES]
    brief = plan.get("rationale", "") or ""
    brand_id = getattr(pack, "id", None)
    # "images_only" → each slide is the on-brand AI image with ONLY the brand logo
    # composited; no headline/body/CTA copy is drawn. "text" (default) overlays the
    # per-slide copy via the deterministic Stage-3 renderer.
    images_only = plan.get("text_mode") == "images_only"
    try:
        refs = rl.reference_images_for(brand_id, "carousel", brief=brief, k=2)
        logo = pipeline.brand_logo_png(brand_id)
    except Exception as exc:  # noqa: BLE001 - setup failed → all-text fallback
        logger.warning("carousel setup failed (%s); using text-frame fallback", exc)
        return _emit(on_artifact, [_carousel_text_fallback(pack, fr, fr.get("index", i + 1),
                                                           images_only)
                                   for i, fr in enumerate(frames)])

    def _render_slide(ordinal: int, fr: dict) -> tuple[int, Artifact]:
        idx = fr.get("index", ordinal + 1)
        role = fr.get("role", "body")
        try:
            run = pipeline.establish_base(
                brand_id, "1:1", reference_images=refs,
                subject=fr.get("subject") or None,
            )
            if images_only:
                # No copy — composite just the logo onto the approved base image.
                # ``subheadings=[]`` clears the run's default sub-text so nothing but
                # the logo is drawn (headline/cta are already empty here).
                png = pipeline.render_frame_on_base(run, logo_png=logo, subheadings=[])
                return idx, (f"frame-{idx:02d}.png", png, "image/png")
            # Layout brain: look at THIS slide's image and place the text in the
            # clean negative space away from the subject (font + gradient stay locked).
            # Carousel text is pinned to a left/right side column (``sides_only``) so a
            # multi-line headline + body can never overflow downward into the subject's
            # face/body — the failure mode of a full-width top/bottom band.
            base_png = pipeline.approved_base_png(run)
            layout = (
                layout_brain.decide_placement(
                    base_png,
                    headline=fr.get("headline", ""),
                    body=fr.get("body", ""),
                    has_cta=(role == "cta"),
                    sides_only=True,
                )
                if base_png else None
            )
            png = pipeline.render_frame_on_base(
                run,
                headline=fr.get("headline", ""),
                highlight=fr.get("highlight", ""),
                subheadings=[fr.get("body", "")] if fr.get("body") else None,
                cta=fr.get("headline", "") if role == "cta" else "",
                logo_png=logo,
                layout=layout,
            )
            return idx, (f"frame-{idx:02d}.png", png, "image/png")
        except Exception as exc:  # noqa: BLE001 - one slide failing ≠ whole set
            logger.warning("slide %s generation failed (%s); text-frame fallback", idx, exc)
            return idx, _carousel_text_fallback(pack, fr, idx, images_only)

    workers = max(1, min(_CAROUSEL_CONCURRENCY, len(frames)))
    results: list[tuple[int, Artifact]] = []
    with concurrent.futures.ThreadPoolExecutor(max_workers=workers) as ex:
        futures = [ex.submit(_render_slide, i, fr) for i, fr in enumerate(frames)]
        for fut in concurrent.futures.as_completed(futures):
            idx, art = fut.result()
            results.append((idx, art))
            if on_artifact:  # emit each slide the moment it finishes (live progress)
                try:
                    on_artifact(art)
                except Exception:  # noqa: BLE001 - progress is best-effort
                    logger.warning("on_artifact callback failed for slide %s", idx, exc_info=True)
    # Frames complete out of order; sort by real frame index for the returned set.
    results.sort(key=lambda t: t[0])
    return [art for _idx, art in results]


def build_blog_images(plan: dict[str, Any], pack: Any,
                      on_artifact: OnArtifact = None) -> list[Artifact]:
    """Real, photorealistic blog imagery — one AI-generated image per visual in
    the plan (cover + in-article), straight from the plan's rich ``visual``
    prompts. A blog hero/section image is just a clean editorial photo: NO brand
    gradient, NO baked-in text — exactly what a good single prompt produces.

    Falls back to a branded gradient frame PER image only if generation is
    unavailable (no provider / API error), so the user always gets a downloadable
    asset instead of nothing.
    """
    from .. import providers, reference_library as rl

    spec = rl.type_spec("blog")
    w, h = spec.get("target_dims", (1600, 900))
    ar = spec.get("aspect_ratio", "16:9")
    try:
        provider = providers.get_provider(agent_id="a1")  # the GD agent's image model
    except Exception as exc:  # noqa: BLE001 - no provider configured → all fallback
        logger.warning("blog: no image provider (%s); using gradient frames", exc)
        provider = None

    def _image(prompt: str, fb_head: str, fb_body: str, role: str) -> bytes:
        """A real generated image for ``prompt``; a branded gradient frame if that
        isn't possible (so one failure never sinks the set)."""
        if provider is not None and (prompt or "").strip():
            try:
                png, _mime = provider.generate(
                    prompt, width=w, height=h, aspect_ratio=ar, image_size="2K",
                )
                return png
            except Exception as exc:  # noqa: BLE001 - one image failing ≠ whole set
                logger.warning("blog image generation failed (%s); gradient fallback", exc)
        return _render_text_frame(pack, (w, h), headline=fb_head, body=fb_body,
                                  role=role, footer=getattr(pack, "name", ""))

    out: list[Artifact] = []
    cover = plan.get("cover", {})
    out.append(("cover.png",
                _image(cover.get("visual", ""), cover.get("title", "Blog"),
                       cover.get("subtitle", ""), "cover"),
                "image/png"))
    for i, inl in enumerate(plan.get("inline", []), 1):
        out.append((f"inline-{i:02d}.png",
                    _image(inl.get("visual", ""), inl.get("caption", f"Figure {i}"),
                           inl.get("visual", ""), "body"),
                    "image/png"))
    return _emit(on_artifact, out)


# --------------------------------------------------------------------------- #
# Brochure (reportlab) + Presentation (python-pptx) — lazy / optional engines
# --------------------------------------------------------------------------- #

def _reportlab_available() -> bool:
    try:
        import reportlab  # noqa: F401
        return True
    except Exception:
        return False


def _pptx_available() -> bool:
    try:
        import pptx  # noqa: F401
        return True
    except Exception:
        return False


def engine_status() -> dict[str, bool]:
    """Which output engines are installed on this host."""
    return {
        "carousel": True,   # Pillow
        "blog": True,       # Pillow
        "brochure": _reportlab_available(),
        "presentation": _pptx_available(),
    }


# A4-ish portrait page for the generated brochure spreads, and how many pages we
# parallelise at once (each page is its own image generation).
_BROCHURE_AR = "4:5"
_BROCHURE_CONCURRENCY = 4


def _study_brand_brochures(brand_id: Optional[str]) -> str:
    """Vision pass over the brand's REAL brochures (the Drive-synced references).

    Returns a tight 'house style' summary — how the subject blends into the
    background, how the brand gradient flows, where photos sit vs text — used to
    steer the new on-brand page generation. ``""`` when no key/refs/vision."""
    try:
        from app.services.openrouter import analyze_images
        from .. import reference_library as rl
    except Exception:
        return ""
    try:
        refs = rl.reference_images_for(brand_id, "brochure", k=3)
    except Exception:  # noqa: BLE001
        refs = []
    if not refs:
        return ""
    prompt = (
        "You are a brand art director. Study these existing brand brochures and "
        "describe, in 2-3 tight sentences, the visual SYSTEM to reuse on new pages: "
        "how the subject/photo blends into the background, how the brand gradient "
        "flows, and where photos sit versus text. Be concrete and brief."
    )
    try:
        notes = analyze_images(prompt, refs[:3])
        return " ".join((notes or "").split())[:500]
    except Exception:  # noqa: BLE001 - study is additive; never block generation
        logger.warning("brochure design-study failed; generating without it", exc_info=True)
        return ""


def build_brochure_pdf(plan: dict[str, Any], pack: Any,
                       on_artifact: OnArtifact = None) -> list[Artifact]:
    """Designed brochure: cards on a calm brand background, composed on a grid and
    assembled into a selectable-text PDF. Fully deterministic (no image provider),
    so it never silently degrades to plain text."""
    if not _reportlab_available():
        raise RuntimeError(
            "Brochure (PDF) generation needs reportlab — add it to requirements "
            "and `pip install reportlab`."
        )
    artifact, _rasters = _designed_brochure_pdf(plan, pack)
    return _emit(on_artifact, [artifact])


def _designed_brochure_pdf(plan: dict[str, Any], pack: Any, *,
                           logo_png: Optional[bytes] = None) -> tuple[Artifact, list[bytes]]:
    """Compose the plan into templated pages, render each to a PNG via the brochure
    layout engine, and assemble the hybrid (image + invisible selectable text) PDF.
    Returns the PDF artifact and the per-page rasters (for page caching)."""
    palette = _brand_palette(pack)
    pages = brochure_layout.compose_brochure(plan)

    def _fonts(size: int, name: Optional[str] = None):
        return _load_font(pack, size, name)

    rasters: list[bytes] = []
    assembled: list[tuple] = []
    for page in pages:
        png = brochure_layout.render_page(
            page, size=brochure_render._BROCHURE_PAGE, palette=palette,
            font_loader=_fonts, logo_png=logo_png,
        )
        rasters.append(png)
        assembled.append((png, {}, page.get("text_lines", [])))

    pdf = _assemble_brochure_pdf(assembled)
    cover = plan.get("cover", {})
    fname = _slug(cover.get("title", "brochure")) + ".pdf"
    return (fname, pdf, "application/pdf"), rasters


def _pdf_wrap(text: str, width: int) -> list[str]:
    words = (text or "").split()
    lines, cur = [], ""
    for w in words:
        if len(cur) + len(w) + 1 <= width or not cur:
            cur = f"{cur} {w}".strip()
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines or [""]


# Each brochure page is a full photographic 2K–4K composite. Embedding those as
# lossless PNG produces tens-of-MB PDFs (slow/failed downloads). Downscale to a
# crisp print resolution and embed as JPEG — a ~10× size cut with no visible loss
# at brochure scale.
_PDF_PAGE_MAX_W = 1654   # ~200 DPI at A4 width
_PDF_JPEG_QUALITY = 85


def _compact_page(png_bytes: bytes) -> bytes:
    """Downscale a page raster to print resolution and re-encode as JPEG."""
    im = Image.open(io.BytesIO(png_bytes)).convert("RGB")
    if im.width > _PDF_PAGE_MAX_W:
        im = im.resize((_PDF_PAGE_MAX_W, round(_PDF_PAGE_MAX_W * im.height / im.width)), Image.LANCZOS)
    out = io.BytesIO()
    im.save(out, format="JPEG", quality=_PDF_JPEG_QUALITY, optimize=True)
    return out.getvalue()


def _assemble_brochure_pdf(pages: list[tuple]) -> bytes:
    """Place each rendered page image full-bleed, then write the page's text as an
    INVISIBLE layer (PDF text render mode 3) so the document stays selectable and
    searchable even though the visible text is the crisp brand-font raster."""
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    page_w = 595.0  # ~A4 width in points
    first = Image.open(io.BytesIO(pages[0][0]))
    iw, ih = first.size
    page_h = page_w * (ih / iw)

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=(page_w, page_h))
    for png, layout, lines in pages:
        c.drawImage(ImageReader(io.BytesIO(_compact_page(png))), 0, 0, width=page_w, height=page_h,
                    preserveAspectRatio=False)
        # Invisible, selectable text roughly in the same zone as the visible copy.
        zone = (layout or {}).get("placement", "left")
        x = page_w * (0.55 if zone == "right" else 0.08)
        text = c.beginText(x, page_h * 0.9)
        text.setTextRenderMode(3)  # invisible — visual text comes from the image
        text.setFont("Helvetica", 11)
        for ln in lines:
            for wrapped in _pdf_wrap(ln, 56):
                text.textLine(wrapped)
        c.drawText(text)
        c.showPage()
    c.save()
    return buf.getvalue()


def build_presentation_pptx(plan: dict[str, Any], pack: Any,
                            on_artifact: OnArtifact = None) -> list[Artifact]:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError(
            "Presentation (PPTX) generation needs python-pptx — add it to "
            "requirements and `pip install python-pptx`."
        ) from exc

    pal = _brand_palette(pack)
    deep = RGBColor(*_hex_to_rgb(pal["deep"]))
    accent = RGBColor(*_hex_to_rgb(pal["accent"]))
    ink = RGBColor(*_hex_to_rgb(pal["text"]))
    font_name = getattr(pack, "font_family", None) or "Calibri"

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _band(slide) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.35))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

    slides = plan.get("slides", [])
    for i, sl in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        is_title = i == 0
        _band(slide)
        # Title
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.9 if is_title else 0.7),
                                      Inches(11.7), Inches(2.0 if is_title else 1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = sl.get("title", "")
        p = tf.paragraphs[0]
        p.font.size = Pt(40 if is_title else 30)
        p.font.bold = True
        p.font.name = font_name
        p.font.color.rgb = deep
        # Bullets
        bullets = sl.get("bullets") or []
        if bullets:
            bb = slide.shapes.add_textbox(Inches(0.9), Inches(3.0 if is_title else 2.2),
                                          Inches(11.5), Inches(4.0))
            btf = bb.text_frame
            btf.word_wrap = True
            for j, b in enumerate(bullets):
                para = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
                para.text = f"•  {b}"
                para.font.size = Pt(20 if is_title else 18)
                para.font.name = font_name
                para.font.color.rgb = ink if not is_title else accent
                para.space_after = Pt(10)
        notes = sl.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buf = io.BytesIO()
    prs.save(buf)
    title = slides[0].get("title", "presentation") if slides else "presentation"
    return _emit(on_artifact, [(_slug(title) + ".pptx",
                                buf.getvalue(),
                                "application/vnd.openxmlformats-officedocument.presentationml.presentation")])


# --------------------------------------------------------------------------- #
# Dispatch + packaging
# --------------------------------------------------------------------------- #

_BUILDERS = {
    "carousel": build_carousel_frames,
    "blog": build_blog_images,
    "brochure": build_brochure_pdf,
    "presentation": build_presentation_pptx,
}


def build(creative_type: str, plan: dict[str, Any], pack: Any,
          on_artifact: OnArtifact = None) -> list[Artifact]:
    """Render a reviewed plan into one or more downloadable artifacts.

    ``on_artifact`` (optional) is called for each artifact the moment it is ready,
    so the caller can persist + show it while the rest of the set is still rendering.
    """
    builder = _BUILDERS.get(creative_type)
    if not builder:
        raise ValueError(f"No document builder for creative type: {creative_type}")
    return builder(plan, pack, on_artifact=on_artifact)


def _slug(text: str) -> str:
    import re
    s = re.sub(r"[^a-z0-9]+", "-", (text or "").lower()).strip("-")
    return s or "creative"


def zip_artifacts(artifacts: list[Artifact], archive_name: str = "creative") -> Artifact:
    """Bundle multiple artifacts (e.g. carousel frames) into one .zip artifact."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, data, _mime in artifacts:
            zf.writestr(name, data)
    return (f"{_slug(archive_name)}.zip", buf.getvalue(), "application/zip")
